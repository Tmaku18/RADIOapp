#!/usr/bin/env python3
"""
Generate professional Word (.docx) and PowerPoint (.pptx) documents
from the markdown analysis for client presentation.
"""

from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.enum.style import WD_BUILTIN_STYLE
import datetime
from pptx import Presentation
from pptx.util import Inches as PptInches
from pptx.enum.text import PP_ALIGN
import re
import os

def create_word_document():
    """Create a professional Word report."""
    doc = Document()
    
    # Set up document
    doc.core_properties.title = "RadioApp Discover Me Development Journey"
    doc.core_properties.author = "Tanaka Makuvaza"
    doc.core_properties.created = datetime.datetime.now()
    
    # Add title
    title = doc.add_heading('RadioApp / Discover Me Development Journey', 0)
    title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    
    subtitle = doc.add_heading('Client Presentation Report - May 2026', 1)
    subtitle.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    
    doc.add_paragraph('Prepared by Tanaka Makuvaza, M.S. Computer Science')
    doc.add_paragraph('Georgia State University')
    doc.add_paragraph(f'Generated on: {datetime.datetime.now().strftime("%B %d, %Y")}')
    
    doc.add_paragraph()  # Spacer
    
    # Executive Summary from the report
    doc.add_heading('Executive Summary', 1)
    summary_text = (
        "Over the past five months, we have built a sophisticated platform that started as a simple "
        "underground artist radio station and has evolved into **Discover Me** - a professional "
        "networking hub for music creators, artists, and service providers (\"Catalysts\" - producers, "
        "mix engineers, videographers, designers, and marketers).\n\n"
        "The development followed a disciplined workflow of detailed planning, status reviews, "
        "focused implementation, comprehensive testing, and rapid bug resolution. The result is a "
        "reliable, fair, and scalable system with transparent radio streaming, rich analytics, "
        "professional networking features, and strong business foundations."
    )
    doc.add_paragraph(summary_text)
    
    doc.add_heading('Development Timeline Overview', 1)
    timeline_summary = (
        "• January 2026: MVP Foundation and completion (Status Reviews 1 & 2)\n"
        "• February-March 2026: Feature expansion and major pivot to Discover Me networking\n"
        "• April 2026: Radio reliability crisis and comprehensive fixes (30+ bugs addressed)\n"
        "• May 2026: Analytics maturation, Refinery completion, compliance features\n\n"
        "Key technical improvements included moving from in-memory state to a robust Redis + "
        "Supabase hybrid architecture, adding play decision logging for transparency, atomic credit "
        "operations, and extensive testing procedures."
    )
    doc.add_paragraph(timeline_summary)
    
    doc.add_heading('Key Challenges & Solutions', 1)
    challenges = [
        ("Early MVP Issues", "Radio state loss on restart, incomplete mobile navigation, no automatic credit deduction, songs stuck pending"),
        ("April Scaling Crisis", "Stuttering/skipping playback, Supabase overload (503 errors, high disk IO), listener count errors"),
        ("Solutions Implemented", "Redis hybrid state management with database checkpoints, play_decision_log for transparency, "
                                  "atomic RPCs for credits, server-side FFmpeg duration validation, caching and circuit breakers, "
                                  "comprehensive test matrix (Master E2E + 14 radio logic tests)")
    ]
    
    for title, description in challenges:
        doc.add_heading(title, 2)
        doc.add_paragraph(description)
    
    doc.add_heading('Current Capabilities', 1)
    capabilities = [
        "Reliable continuous radio streaming with realtime chat and smooth transitions",
        "Professional artist analytics dashboard with real per-song statistics",
        "Refinery system for paid submissions and structured reviews",
        "Discover Me networking - profiles, service marketplace, job board, secure messaging",
        "Transparent credit system with audit logging",
        "Cross-platform consistency (web, mobile, backend)",
        "Self-service account deletion for compliance"
    ]
    
    for cap in capabilities:
        doc.add_paragraph(cap, style='List Bullet')
    
    doc.add_heading('Development Workflow', 1)
    workflow = (
        "1. Detailed planning documents (.cursor/plans/ with 27+ files including Status Reviews)\n"
        "2. Regular status reviews to track progress and gaps\n"
        "3. Iterative implementation with clear git commits (497 total)\n"
        "4. Comprehensive testing per Testing_Procedures.md (E2E flows, radio logic tests R1-R14)\n"
        "5. Systematic bug fixing with architectural improvements\n\n"
        "This methodical approach ensured reliability and transparency throughout development."
    )
    doc.add_paragraph(workflow)
    
    doc.add_heading('Next Steps & Recommendations', 1)
    next_steps = [
        "Complete full Discover Me browse and messaging UI",
        "Implement competitions, spotlight features, and livestream capabilities",
        "Finalize comprehensive automated test suite",
        "Prepare for broader user acquisition and marketing",
        "Schedule client demo of current analytics dashboard and radio player"
    ]
    
    for step in next_steps:
        doc.add_paragraph(step, style='List Bullet')
    
    doc.add_paragraph()
    doc.add_paragraph('This document was auto-generated from git history, all planning documents, '
                     'status reviews, and current codebase analysis.')
    
    # Save
    output_path = 'presentation/RadioApp_Development_Journey_Report.docx'
    doc.save(output_path)
    print(f"Word document created: {output_path}")
    return output_path


def create_powerpoint():
    """Create a professional PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = PptInches(13.33)
    prs.slide_height = PptInches(7.5)
    
    # Set theme colors (professional blue/green)
    color_scheme = prs.slide_master.slide_layouts[0].background.fill.fore_color
    # Note: In real use we'd customize more, but this creates clean slides
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Discover Me"
    subtitle.text = "From Radio Streaming MVP to Professional Creator Platform\nDevelopment Journey & Analysis\nMay 2026"
    
    # Add presenter info
    left = PptInches(1)
    top = PptInches(5)
    width = PptInches(11)
    height = PptInches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    frame = txBox.text_frame
    frame.text = "Tanaka Makuvaza, M.S. Computer Science\nGeorgia State University"
    
    # Slide 2: Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Agenda"
    content = slide.placeholders[1]
    content.text = ("• Project Overview & Vision\n"
                   "• Development Workflow & Methodology\n"
                   "• Complete Timeline (Jan-May 2026)\n"
                   "• MVP Completion & Major Pivot\n"
                   "• The April Reliability Crisis & Fixes\n"
                   "• Current Platform Capabilities\n"
                   "• Testing & Quality Approach\n"
                   "• Recent Analytics & Refinery Enhancements\n"
                   "• Next Steps & Roadmap\n"
                   "• Summary and Q&A")
    
    # Add more slides based on the outline (abbreviated for key content)
    slides_data = [
        ("What is Discover Me?", 
         "Professional networking platform for music creators\n\n"
         "• Reliable continuous radio streaming\n"
         "• Artist analytics and Refinery (paid reviews)\n"
         "• Profiles, portfolios, job board for Catalysts\n"
         "• Credit system with full transparency logging\n\n"
         "Radio = powerful engagement engine\n"
         "Networking = core business value", 
         "We evolved from pure radio pay-per-play to a complete creator economy platform."),
        
        ("Development Workflow", 
         "1. Detailed Cursor plans (27+ documents)\n"
         "2. Status Reviews as progress checkpoints\n"
         "3. Focused implementation sprints with parity\n"
         "4. Comprehensive Testing_Procedures.md\n"
         "5. Systematic bug fixing with architecture upgrades\n\n"
         "497 git commits with clear descriptions", 
         "Methodical, transparent, and quality-focused approach using plans, reviews, and rigorous testing."),
        
        ("Timeline Overview", 
         "Jan 2026: MVP Foundation & Completion (Status Reviews 1 & 2)\n"
         "Feb-Mar 2026: Feature expansion + Discover Me pivot\n"
         "April 2026: Radio regression & 30+ bug fixes\n"
         "May 2026: Analytics, Refinery completion, compliance\n\n"
         "See timeline.md for full Mermaid graphic", 
         "Clear progression from concept to production platform with major reliability improvements in April."),
        
        ("April Reliability Crisis", 
         "Problems at scale:\n"
         "• Radio stuttering/skipping/restarting\n"
         "• Supabase 503 errors and high disk IO\n"
         "• Listener count inaccuracies\n\n"
         "Root cause: Initial in-memory state management", 
         "These issues appeared as usage grew. We responded with a comprehensive 2-week architectural overhaul."),
        
        ("Solutions Implemented", 
         "• Redis + Supabase hybrid state (stateless backend)\n"
         "• play_decision_log table for full transparency\n"
         "• Atomic credit RPCs and server-side duration validation\n"
         "• Caching, circuit breakers, fallbacks, parallel queries\n"
         "• Comprehensive test matrix (E2E + 14 radio tests)", 
         "Result: Stable, trustworthy, and scalable radio streaming. Platform is now production-ready."),
        
        ("Current Capabilities", 
         "• Reliable radio with realtime chat & smooth playback\n"
         "• Rich artist statistics dashboard (new DB tables)\n"
         "• Professional Refinery for paid reviews\n"
         "• Discover Me networking (profiles, jobs, marketplace)\n"
         "• Transparent credit system with audit trails\n"
         "• Cross-platform (Web + Mobile parity)", 
         "The platform today delivers real business value with strong technical foundations."),
        
        ("Testing & Quality", 
         "Testing_Procedures.md serves as living specification:\n"
         "• All backend modules implemented and documented\n"
         "• Master 10-step end-to-end user journey\n"
         "• 14 specific Radio Logic test cases (R1-R14)\n"
         "• Unit, integration, load, and security testing framework", 
         "Quality was built-in from the start. Prevents regression as new features are added."),
        
        ("Recent Work (May 2026)", 
         "• New artist stats page and analytics.service.ts\n"
         "• Database migrations for song stats and daily metrics\n"
         "• Complete Refinery overhaul with structured reviews\n"
         "• Self-service account deletion compliance page\n"
         "• Enhanced My Songs with real performance data", 
         "Development momentum remains strong with focus on analytics and creator tools."),
        
        ("Next Steps", 
         "• Complete Discover Me UI (browse feed, messaging, job flows)\n"
         "• Add competitions, spotlight features, and livestreams\n"
         "• Execute full automated test suite\n"
         "• Polish interfaces and prepare for launch\n"
         "• Schedule client demo of current platform", 
         "The foundation is solid. We are well-positioned for growth."),
        
        ("Summary", 
         "• Methodical development from MVP to full platform in 5 months\n"
         "• Successfully overcame scaling challenges with architectural improvements\n"
         "• Built transparent, reliable, and scalable system\n"
         "• Strong combination of radio engagement + professional networking\n"
         "• Production-ready with clear roadmap ahead", 
         "Thank You\n\nQ&A\n\nReady for live demo of radio player, analytics dashboard, and admin tools.")
    ]
    
    for slide_title, content_text, notes_text in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout for more control
        title_shape = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(12.3), PptInches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = slide_title
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
        
        # Add content
        content_shape = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.5), PptInches(11.7), PptInches(4.5))
        content_frame = content_shape.text_frame
        content_frame.text = content_text
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.add_paragraph().text = notes_text
    
    # Save
    output_path = 'presentation/RadioApp_Development_Presentation.pptx'
    prs.save(output_path)
    print(f"PowerPoint presentation created: {output_path}")
    return output_path


if __name__ == "__main__":
    print("Generating client Word and PowerPoint documents...")
    word_file = create_word_document()
    ppt_file = create_powerpoint()
    print("\n✅ Successfully generated:")
    print(f"   • {word_file}")
    print(f"   • {ppt_file}")
    print("\nThese are professional, ready-to-use Office documents for clients.")
    print("Open them in Microsoft Word and PowerPoint to review and present.")
